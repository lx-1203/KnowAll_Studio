"""Document parsing engine - local file parsing, no API calls"""
import hashlib
import logging
import re
from pathlib import Path
from dataclasses import dataclass, field
from app.config import settings

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    """Result of document parsing"""
    text: str
    page_count: int
    metadata: dict = field(default_factory=dict)
    headings: list = field(default_factory=list)   # list[HeadingNode] from docling
    images: list = field(default_factory=list)     # list[ImageItem] from docling
    tables: list = field(default_factory=list)     # list[TableItem] from docling


@dataclass
class TextChunk:
    """A single text chunk after splitting"""
    index: int
    text: str
    token_count: int
    content_hash: str
    page_range: str = ""


class DocumentParser:
    """Parse various document formats to plain text (100% local)."""

    async def parse(self, file_path: str, file_type: str) -> ParsedDocument:
        # Try Docling first for formats it excels at
        if settings.use_docling and file_type in ("pdf", "docx", "pptx", "html"):
            try:
                return await self._parse_with_docling(file_path, file_type)
            except Exception as e:
                logger.warning("Docling parse failed for %s (%s), falling back: %s",
                               file_type, file_path, e)

        parsers = {
            "pdf": self._parse_pdf,
            "docx": self._parse_docx,
            "pptx": self._parse_pptx,
            "md": self._parse_markdown,
            "markdown": self._parse_markdown,
            "txt": self._parse_text,
            "text": self._parse_text,
            "png": self._parse_image,
            "jpg": self._parse_image,
            "jpeg": self._parse_image,
            "gif": self._parse_image,
            "bmp": self._parse_image,
            "webp": self._parse_image,
            "xlsx": self._parse_xlsx,
            "csv": self._parse_csv,
            "xmind": self._parse_xmind,
            "url": self._parse_url,
            "html": self._parse_html,
            "py": self._parse_code,
            "js": self._parse_code,
            "ts": self._parse_code,
            "jsx": self._parse_code,
            "tsx": self._parse_code,
            "java": self._parse_code,
            "cpp": self._parse_code,
            "c": self._parse_code,
            "h": self._parse_code,
            "go": self._parse_code,
            "rs": self._parse_code,
            "sql": self._parse_code,
            "yaml": self._parse_code,
            "yml": self._parse_code,
            "json": self._parse_code,
            "xml": self._parse_code,
            "css": self._parse_code,
            "html_code": self._parse_code,
        }
        parser = parsers.get(file_type)
        if not parser:
            raise ValueError(f"Unsupported file type: {file_type}. Supported: {list(parsers.keys())}")
        return await parser(file_path)

    async def _parse_pdf(self, file_path: str) -> ParsedDocument:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        texts = []
        for page in doc:
            texts.append(page.get_text())
        doc.close()
        return ParsedDocument(
            text="\n\n".join(texts),
            page_count=len(texts),
            metadata={"parser": "PyMuPDF"},
        )

    async def _parse_docx(self, file_path: str) -> ParsedDocument:
        import docx
        doc = docx.Document(file_path)
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        return ParsedDocument(
            text="\n\n".join(texts),
            page_count=1,
            metadata={"parser": "python-docx"},
        )

    async def _parse_pptx(self, file_path: str) -> ParsedDocument:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
            if slide_texts:
                texts.append("\n".join(slide_texts))
        return ParsedDocument(
            text="\n\n---\n\n".join(texts),
            page_count=len(texts),
            metadata={"parser": "python-pptx"},
        )

    async def _parse_markdown(self, file_path: str) -> ParsedDocument:
        text = Path(file_path).read_text(encoding="utf-8")
        return ParsedDocument(
            text=text,
            page_count=1,
            metadata={"parser": "builtin"},
        )

    async def _parse_text(self, file_path: str) -> ParsedDocument:
        text = Path(file_path).read_text(encoding="utf-8")
        return ParsedDocument(
            text=text,
            page_count=1,
            metadata={"parser": "builtin"},
        )

    async def _parse_image(self, file_path: str) -> ParsedDocument:
        """OCR image to text using Tesseract or PaddleOCR."""
        text = ""
        parser_used = "none"

        # Try PaddleOCR first (better for Chinese)
        try:
            from paddleocr import PaddleOCR
            ocr = PaddleOCR(lang="ch", show_log=False)
            result = ocr.ocr(file_path)
            if result and result[0]:
                text = "\n".join(line[1][0] for line in result[0] if line)
            parser_used = "PaddleOCR"
        except ImportError:
            pass

        # Fallback to Tesseract
        if not text:
            try:
                import pytesseract
                from PIL import Image
                img = Image.open(file_path)
                text = pytesseract.image_to_string(img, lang="chi_sim+eng")
                parser_used = "Tesseract"
            except ImportError:
                pass

        if not text.strip():
            raise ValueError(
                "OCR failed: could not extract text from image. "
                "Please install PaddleOCR or Tesseract."
            )

        return ParsedDocument(
            text=text,
            page_count=1,
            metadata={"parser": parser_used, "source": "ocr"},
        )

    async def _parse_xmind(self, file_path: str) -> ParsedDocument:
        """Parse XMind mind-map file."""
        try:
            import xmindparser
            # xmindparser supports both dict and list output
            content = xmindparser.xmind_to_dict(file_path)
            texts = []

            def extract_topics(data, depth=0):
                if isinstance(data, list):
                    for item in data:
                        extract_topics(item, depth)
                elif isinstance(data, dict):
                    title = data.get("title", "")
                    if title:
                        texts.append("#" * min(depth + 1, 6) + " " + title)
                    children = data.get("topics", [])
                    if isinstance(children, list):
                        for child in children:
                            extract_topics(child, depth + 1)

            # Handle different xmindparser output formats
            if isinstance(content, list):
                for item in content:
                    extract_topics(item)
            elif isinstance(content, dict):
                sheet = content.get("sheet") or content
                root_topic = sheet.get("rootTopic") or sheet.get("topic", {})
                extract_topics(root_topic)
            else:
                extract_topics(content)

            text = "\n\n".join(texts) if texts else str(content)
            return ParsedDocument(text=text, page_count=1, metadata={"parser": "xmindparser"})
        except ImportError:
            raise ValueError("XMind parsing requires: pip install xmindparser")

    async def _parse_code(self, file_path: str) -> ParsedDocument:
        """Parse source code files with language tagging."""
        ext = Path(file_path).suffix.lstrip(".")
        lang_map = {
            "py": "Python", "js": "JavaScript", "ts": "TypeScript",
            "jsx": "React JSX", "tsx": "React TSX", "java": "Java",
            "cpp": "C++", "c": "C", "h": "C/C++ Header", "go": "Go",
            "rs": "Rust", "sql": "SQL", "yaml": "YAML", "yml": "YAML",
            "json": "JSON", "xml": "XML", "css": "CSS",
        }
        lang = lang_map.get(ext, ext.upper())
        text = Path(file_path).read_text(encoding="utf-8", errors="replace")
        # Wrap code with language tag for better LLM understanding
        text = f"以下为{lang}源代码文件:\n\n```{ext}\n{text}\n```"
        return ParsedDocument(text=text, page_count=1, metadata={"parser": "code", "language": lang})

    async def _parse_url(self, file_path: str) -> ParsedDocument:
        """Fetch and parse a URL. file_path is the URL string."""
        url = file_path
        try:
            import httpx
            # Synchronous fetch for simplicity (runs in async context)
            import asyncio
            async with httpx.AsyncClient(timeout=30) as client:
                response = await client.get(url, follow_redirects=True)
                response.raise_for_status()
                html = response.text
        except ImportError:
            raise ValueError("URL parsing requires: pip install httpx")

        # Simple HTML to text extraction
        import re
        # Remove scripts and styles
        html = re.sub(r"<script[^>]*>[\s\S]*?</script>", "", html, flags=re.IGNORECASE)
        html = re.sub(r"<style[^>]*>[\s\S]*?</style>", "", html, flags=re.IGNORECASE)
        # Remove HTML tags
        text = re.sub(r"<[^>]+>", " ", html)
        # Decode HTML entities
        text = text.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
        text = text.replace("&quot;", "\"").replace("&#39;", "'").replace("&nbsp;", " ")
        # Normalize whitespace
        text = re.sub(r"\s+", " ", text).strip()

        if not text or len(text) < 50:
            raise ValueError("URL returned insufficient content")

        return ParsedDocument(
            text=text, page_count=1,
            metadata={"parser": "url", "source": url[:200]},
        )

    async def _parse_html(self, file_path: str) -> ParsedDocument:
        """Parse a local HTML file."""
        text = Path(file_path).read_text(encoding="utf-8", errors="replace")
        import re
        text = re.sub(r"<script[^>]*>[\s\S]*?</script>", "", text, flags=re.IGNORECASE)
        text = re.sub(r"<style[^>]*>[\s\S]*?</style>", "", text, flags=re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        return ParsedDocument(text=text, page_count=1, metadata={"parser": "html"})

    async def _parse_with_docling(self, file_path: str, file_type: str) -> ParsedDocument:
        """Parse using Docling for structure-aware extraction."""
        from app.core.parsing.docling_parser import docling_parser as dl_parser
        structured = await dl_parser.parse(file_path, file_type)
        return ParsedDocument(
            text=structured.text,
            page_count=structured.page_count,
            metadata=structured.metadata,
            headings=structured.headings,
            images=structured.images,
            tables=structured.tables,
        )


class TextCleaner:
    """Clean extracted text: remove headers, footers, watermarks, ads, redundant whitespace."""

    # Common watermark patterns
    WATERMARK_PATTERNS = [
        r"版权所有.*?翻印必究",
        r"Copyright.*?All rights reserved",
        r"仅供学习交流",
        r"扫描二维码.*",
    ]

    def clean(self, text: str) -> str:
        """Apply all cleaning steps."""
        text = self._remove_watermarks(text)
        text = self._normalize_whitespace(text)
        text = self._remove_redundant_blank_lines(text)
        text = self._remove_short_lines(text, min_chars=3)
        return text.strip()

    def _remove_watermarks(self, text: str) -> str:
        for pattern in self.WATERMARK_PATTERNS:
            text = re.sub(pattern, "", text, flags=re.IGNORECASE)
        return text

    def _normalize_whitespace(self, text: str) -> str:
        # Replace multiple spaces/newlines
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text

    def _remove_redundant_blank_lines(self, text: str) -> str:
        return re.sub(r"\n\s*\n\s*\n", "\n\n", text)

    def _remove_short_lines(self, text: str, min_chars: int = 3) -> str:
        """Remove lines that are too short (page numbers, standalone symbols).

        Preserves heading lines even if they're short (e.g. "pH", "DNA", "1. 概述").
        """
        lines = text.split("\n")
        cleaned = []
        for line in lines:
            stripped = line.strip()
            # Keep if long enough
            if len(stripped) >= min_chars:
                cleaned.append(line)
                continue
            # Keep if it looks like a heading or numbered item
            if re.match(r"^[#]", stripped):     # Markdown heading #
                cleaned.append(line)
                continue
            if re.match(r"^[\d]+[\.\)、]\s", stripped):  # Numbered item (1. 1) 1、)
                cleaned.append(line)
                continue
            if re.match(r"^[IVX]+[\.\)、]", stripped):   # Roman numeral
                cleaned.append(line)
                continue
            if re.match(r"^[一二三四五六七八九十]、", stripped):  # Chinese numbered
                cleaned.append(line)
                continue
            if re.match(r"^（[一二三四五六七八九十\d]）", stripped):  # (一) (1)
                cleaned.append(line)
                continue
            # Drop line - too short and doesn't look like a heading/list item
        return "\n".join(cleaned)


class TextSplitter:
    """Split long documents into chunks suitable for LLM context windows.

    Heading-aware: detects Markdown headings and Chinese-numbered headings
    to avoid splitting heading content from its body.
    """

    # Regex patterns for detecting heading lines
    HEADING_PATTERNS = [
        re.compile(r'^#{1,6}\s+\S'),           # Markdown headings
        re.compile(r'^第[一二三四五六七八九十百千\d]+[章节篇部]'),     # 第X章/节/篇
        re.compile(r'^[\d]+\.[\d]+\s+\S'),     # 1.1 / 2.3 numbered
        re.compile(r'^[一二三四五六七八九十]、'),  # 一、二、
        re.compile(r'^（[一二三四五六七八九十\d]）'),  # (一) (1)
        re.compile(r'^[IVX]+[\.、]\s'),          # I. II. III.
    ]

    def __init__(
        self,
        chunk_size_tokens: int | None = None,
        overlap_ratio: float | None = None,
    ):
        self.chunk_size_tokens = chunk_size_tokens or settings.chunk_size_tokens
        self.overlap_ratio = overlap_ratio or settings.chunk_overlap_ratio

    def _is_heading(self, line: str) -> bool:
        """Check if a line is a heading/section title."""
        stripped = line.strip()
        if not stripped:
            return False
        for pattern in self.HEADING_PATTERNS:
            if pattern.match(stripped):
                return True
        return False

    def split(self, text: str, total_pages: int = 1) -> list[TextChunk]:
        """Split text into overlapping chunks at natural boundaries.

        Heading-aware: detects headings and prefers to start new chunks
        at heading boundaries to preserve the structural relationship
        between headings and their content.
        """
        paragraphs = text.split("\n")
        chunks = []
        current_chunk = []
        current_tokens = 0
        chunk_index = 0
        overlap_size = int(self.chunk_size_tokens * self.overlap_ratio)

        heading_token_buffer = int(self.chunk_size_tokens * 0.15)  # 15% headroom

        for para in paragraphs:
            para_tokens = self._estimate_tokens(para)
            is_heading = self._is_heading(para)

            # Determine if we should start a new chunk
            would_overflow = current_tokens + para_tokens > self.chunk_size_tokens
            near_full = current_tokens > self.chunk_size_tokens - heading_token_buffer

            if would_overflow and current_chunk:
                # Finish current chunk
                chunk_text = "\n".join(current_chunk)
                chunks.append(TextChunk(
                    index=chunk_index,
                    text=chunk_text,
                    token_count=current_tokens,
                    content_hash=self._hash_text(chunk_text),
                    page_range=f"p{max(1, chunk_index * total_pages // max(1, len(paragraphs)))}-p{min(total_pages, (chunk_index+1) * total_pages // max(1, len(paragraphs)))}",
                ))
                chunk_index += 1

                # Start new chunk with overlap
                overlap_text = self._get_overlap(current_chunk, overlap_size)
                current_chunk = overlap_text
                current_tokens = sum(self._estimate_tokens(t) for t in current_chunk)

            # Prefer heading boundary: if near full and new heading, start fresh
            elif is_heading and near_full and len(current_chunk) > 1:
                chunk_text = "\n".join(current_chunk)
                chunks.append(TextChunk(
                    index=chunk_index,
                    text=chunk_text,
                    token_count=current_tokens,
                    content_hash=self._hash_text(chunk_text),
                    page_range=f"p{max(1, chunk_index * total_pages // max(1, len(paragraphs)))}-p{min(total_pages, (chunk_index+1) * total_pages // max(1, len(paragraphs)))}",
                ))
                chunk_index += 1
                current_chunk = []
                current_tokens = 0

            current_chunk.append(para)
            current_tokens += para_tokens

        # Don't forget the last chunk
        if current_chunk:
            chunk_text = "\n".join(current_chunk)
            chunks.append(TextChunk(
                index=chunk_index,
                text=chunk_text,
                token_count=current_tokens,
                content_hash=self._hash_text(chunk_text),
                page_range=f"p{max(1, chunk_index * total_pages // max(1, len(paragraphs)))}-p{total_pages}",
            ))

        return chunks

    def _get_overlap(self, lines: list[str], overlap_tokens: int) -> list[str]:
        """Get the last N lines that fit within overlap token budget."""
        overlap_lines = []
        token_count = 0
        for line in reversed(lines):
            line_tokens = self._estimate_tokens(line)
            if token_count + line_tokens > overlap_tokens:
                break
            overlap_lines.insert(0, line)
            token_count += line_tokens
        return overlap_lines

    def _estimate_tokens(self, text: str) -> int:
        """Rough token estimation: Chinese ~1.5 char/token, English ~4 char/token."""
        chinese_chars = len(re.findall(r"[\u4e00-\u9fff]", text))
        other_chars = len(text) - chinese_chars
        return int(chinese_chars * 0.6 + other_chars * 0.25)

    @staticmethod
    def _hash_text(text: str) -> str:
        normalized = "".join(text.split()).lower()
        return hashlib.sha256(normalized.encode("utf-8")).hexdigest()


# Singleton instances
parser = DocumentParser()
cleaner = TextCleaner()
splitter = TextSplitter()
