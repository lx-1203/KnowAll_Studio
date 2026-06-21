"""PPTX slide-to-image renderer for document preview.
Uses python-pptx + Pillow (already installed) — zero new dependencies.
"""
import os
from io import BytesIO

from PIL import Image, ImageDraw, ImageFont


class PPTXRenderer:
    """Render PPTX slides as PNG images for preview."""

    # Default slide dimensions (16:9) in pixels at 96 DPI
    DEFAULT_WIDTH = 960
    DEFAULT_HEIGHT = 540

    # Try to load a CJK-capable font
    _font_paths = [
        # Windows
        "C:/Windows/Fonts/msyh.ttc",       # Microsoft YaHei
        "C:/Windows/Fonts/simsun.ttc",     # SimSun
        "C:/Windows/Fonts/simhei.ttf",     # SimHei
        # macOS
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        # Linux
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
        "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
    ]

    def _find_font(self, size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
        """Find a suitable font, falling back to default."""
        for path in self._font_paths:
            if os.path.exists(path):
                try:
                    return ImageFont.truetype(path, size)
                except Exception:
                    continue
        return ImageFont.load_default()

    def get_slide_count(self, file_path: str) -> int:
        """Return number of slides in a PPTX file."""
        from pptx import Presentation
        prs = Presentation(file_path)
        return len(prs.slides)

    def render_slide(self, file_path: str, slide_index: int) -> bytes:
        """Render a single slide as PNG, return raw bytes.

        slide_index is 0-based.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor

        prs = Presentation(file_path)
        slides = list(prs.slides)
        if slide_index < 0 or slide_index >= len(slides):
            raise IndexError(f"Slide {slide_index} out of range (0-{len(slides)-1})")

        slide = slides[slide_index]

        # Get slide dimensions
        sw = prs.slide_width or 12192000   # default 16:9 in EMU
        sh = prs.slide_height or 6858000

        # Scale to output image
        scale_x = self.DEFAULT_WIDTH / sw
        scale_y = self.DEFAULT_HEIGHT / sh

        # Create image canvas
        img = Image.new("RGB", (self.DEFAULT_WIDTH, self.DEFAULT_HEIGHT), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        title_font = self._find_font(22)
        body_font = self._find_font(16)
        small_font = self._find_font(12)

        for shape in slide.shapes:
            try:
                left = int(shape.left * scale_x) if shape.left else 0
                top = int(shape.top * scale_y) if shape.top else 0
                width = int(shape.width * scale_x) if shape.width else 100
                height = int(shape.height * scale_y) if shape.height else 50

                if shape.has_text_frame:
                    self._draw_text_frame(draw, shape.text_frame, left, top, width, height,
                                          title_font, body_font, small_font)

                elif shape.shape_type == 13:  # Picture
                    self._draw_picture(img, shape, left, top, width, height)

                elif shape.has_table:
                    self._draw_table(draw, shape.table, left, top, width, height, small_font)

                # Draw shape border for non-text placeholders
                if not shape.has_text_frame and shape.shape_type != 13:
                    draw.rectangle([left, top, left + width, top + height],
                                   outline=(200, 200, 200), width=1)

            except Exception:
                # Skip shapes that fail to render
                pass

        # Add slide number
        draw.text((self.DEFAULT_WIDTH - 60, self.DEFAULT_HEIGHT - 30),
                  f"{slide_index + 1}/{len(slides)}",
                  fill=(180, 180, 180), font=small_font)

        # Encode as PNG bytes
        buf = BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    def _draw_text_frame(self, draw, text_frame, left, top, width, height,
                         title_font, body_font, small_font):
        """Draw text frame content onto the canvas."""
        y = top + 4
        line_spacing = 20

        for para in text_frame.paragraphs:
            if not para.text.strip():
                y += line_spacing // 2
                continue

            # Choose font based on paragraph level/formatting
            try:
                font_size = para.runs[0].font.size if para.runs and para.runs[0].font.size else None
            except Exception:
                font_size = None

            if font_size:
                import pptx.util
                pt_size = font_size / 12700  # EMU to points
                font = self._find_font(int(min(pt_size * 1.2, 24)))
            elif para.level == 0:
                font = title_font
            elif para.level <= 2:
                font = body_font
            else:
                font = small_font

            # Word wrap: break text into lines fitting in the shape width
            text = para.text.strip()
            lines = self._wrap_text(text, font, width - 8)

            for line in lines:
                if y + line_spacing > top + height:
                    break  # text overflow, stop drawing
                draw.text((left + 4, y), line, fill=(50, 50, 50), font=font)
                y += line_spacing

    def _draw_picture(self, img, shape, left, top, width, height):
        """Extract and place an embedded picture on the canvas."""
        try:
            image_blob = shape.image.blob
            content_type = shape.image.content_type
            ext = "png" if "png" in content_type else "jpg"
            pic = Image.open(BytesIO(image_blob))
            # Scale to fit while maintaining aspect ratio
            pic_w, pic_h = pic.size
            scale = min(width / pic_w, height / pic_h)
            new_w = int(pic_w * scale)
            new_h = int(pic_h * scale)
            if new_w > 0 and new_h > 0:
                pic = pic.resize((new_w, new_h), Image.LANCZOS)
                # Center in shape area
                px = left + (width - new_w) // 2
                py = top + (height - new_h) // 2
                if pic.mode == "RGBA":
                    img.paste(pic, (px, py), pic)
                else:
                    img.paste(pic, (px, py))
        except Exception:
            pass

    def _draw_table(self, draw, table, left, top, width, height, font):
        """Draw a basic table grid with text."""
        rows = len(table.rows)
        cols = len(table.columns)
        if rows == 0 or cols == 0:
            return

        cell_w = width // cols
        cell_h = min(height // rows, 30)

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cx = left + c * cell_w
                cy = top + r * cell_h
                # Draw cell border
                draw.rectangle([cx, cy, cx + cell_w, cy + cell_h],
                               outline=(180, 180, 180), width=1)
                # Draw header row with background
                if r == 0:
                    draw.rectangle([cx + 1, cy + 1, cx + cell_w - 1, cy + cell_h - 1],
                                   fill=(230, 230, 240))
                # Draw text (first line only to save space)
                text = cell.text.strip()[:30]
                if text:
                    draw.text((cx + 4, cy + 4), text, fill=(50, 50, 50), font=font)

    def _wrap_text(self, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
        """Simple word wrapping for Pillow text rendering."""
        words = text
        lines = []
        current_line = ""

        for char in words:
            test_line = current_line + char
            try:
                bbox = font.getbbox(test_line) if hasattr(font, 'getbbox') else font.getsize(test_line)
                text_width = bbox[2] if isinstance(bbox, tuple) and len(bbox) == 4 else bbox[0]
            except Exception:
                text_width = len(test_line) * 10

            if text_width > max_width and current_line:
                lines.append(current_line)
                current_line = char
            else:
                current_line = test_line

        if current_line:
            lines.append(current_line)

        return lines or [text]


# Singleton
pptx_renderer = PPTXRenderer()
